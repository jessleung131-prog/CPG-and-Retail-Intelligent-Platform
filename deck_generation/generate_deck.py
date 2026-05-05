"""
Executive PPTX deck generator.

Builds a clean, business-ready presentation from the insights report.

Sections:
  1. Title slide
  2. Executive Summary
  3. Marketing Channel Highlights
  4. Budget Optimizer
  5. Sales & Funnel Highlights
  6. Risks & Opportunities
  7. Recommended Actions

Run:
    python deck_generation/generate_deck.py
Output:
    outputs/executive_deck.pptx
"""
from __future__ import annotations

import json
import os
import sys
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "outputs"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
TEAL   = RGBColor(0x00, 0x9B, 0x8E)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x6B, 0x7A, 0x90)
RED    = RGBColor(0xF0, 0x44, 0x38)
AMBER  = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x12, 0xB7, 0x6A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Slide builder helpers ─────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide, left, top, width, height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = NAVY,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox_multiline(slide, left, top, width, height, lines: list, align=PP_ALIGN.LEFT):
    """Add textbox with multiple paragraphs with independent formatting.

    lines: list of (text, size, bold, color) tuples
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _bullet_list(slide, left, top, width, height, items: list[str], font_size=16, color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def _kpi_box(slide, left, top, label: str, value: str, color: RGBColor = TEAL,
             val_y_offset=Inches(0.15), lbl_y_offset=Inches(0.85), box_w=Inches(4.2), box_h=Inches(1.3)):
    _add_rect(slide, left, top, box_w, box_h, color)
    _add_textbox(slide, left + Inches(0.1), top + val_y_offset, box_w - Inches(0.2),
                 Inches(0.5), value, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, left + Inches(0.1), top + lbl_y_offset, box_w - Inches(0.2),
                 Inches(0.4), label, font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ── Section slide builders ────────────────────────────────────────────────────

def _slide_title(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, NAVY)

    # Author top-right
    _add_textbox(slide, Inches(11.5), Inches(0.2), Inches(1.6), Inches(0.4),
                 "Jess L.", font_size=11, color=WHITE, align=PP_ALIGN.RIGHT)

    # Title
    _add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.0),
                 "CPG & Retail Intelligence Platform",
                 font_size=40, bold=True, color=WHITE)

    # Subtitle
    _add_textbox(slide, Inches(1), Inches(2.4), Inches(8), Inches(0.6),
                 "Executive Performance Report",
                 font_size=22, color=TEAL)

    # Period
    _add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.45),
                 "2023–2024  ·  Full Year Analysis",
                 font_size=14, color=GREY)

    # Teal bottom bar
    _add_rect(slide, 0, Inches(5.7), SLIDE_W, Inches(1.4), TEAL)

    # KPI chips inside teal bar
    kpi_chips = [
        ("$78.8M Revenue",       Inches(0.8)),
        ("$15.0M Media Spend",   Inches(4.0)),
        ("2.21× Blended ROAS", Inches(7.2)),
        ("2.6% Win Rate",        Inches(10.5)),
    ]
    for value_label, x in kpi_chips:
        # Split into value and label parts
        parts = value_label.split(" ", 1)
        value_part = parts[0]
        label_part = parts[1] if len(parts) > 1 else ""
        _add_textbox(slide, x, Inches(5.8), Inches(2.8), Inches(0.5),
                     value_part, font_size=22, bold=True, color=WHITE)
        _add_textbox(slide, x, Inches(5.8) + Inches(0.55), Inches(2.8), Inches(0.4),
                     label_part, font_size=10, color=WHITE)

    # Footer
    _add_textbox(slide, Inches(1), Inches(7.15), Inches(8), Inches(0.3),
                 "Generated May 04, 2026  |  Confidential",
                 font_size=10, color=GREY)


def _slide_exec_summary(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    # Navy header
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 "Executive Summary", font_size=24, bold=True, color=WHITE)

    # Left column — 3 fixed bullet points
    bullets = [
        "$78.8M total revenue — online +8.3% MoM, offline dominates at 63% of total",
        "Blended ROAS of 2.21× signals media inefficiency — Email (6.39×) is severely underfunded vs. Display (1.81×) and Reddit (1.61×)",
        "2.62% funnel win rate is the most urgent issue — 97.4% of 25K leads never convert",
    ]
    _bullet_list(slide, Inches(0.6), Inches(1.1), Inches(7.5), Inches(4.5),
                 bullets, font_size=13, color=NAVY)

    # Right column — 4 KPI cards
    cards = [
        (NAVY,  "$78.8M",  "Combined Revenue",      Inches(1.0)),
        (TEAL,  "$15.0M",  "Media Spend",            Inches(2.45)),
        (AMBER, "2.21×",   "Blended ROAS",           Inches(3.9)),
        (RED,   "2.6%",    "Win Rate (target 5–8%)", Inches(5.35)),
    ]
    for color, value, label, y in cards:
        _kpi_box(slide, Inches(8.6), y, label, value, color=color,
                 val_y_offset=Inches(0.15), lbl_y_offset=Inches(0.85),
                 box_w=Inches(4.2), box_h=Inches(1.3))


def _slide_channel_highlights(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    # Navy header
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 "Marketing Channel Performance", font_size=24, bold=True, color=WHITE)

    # Section labels
    _add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.35),
                 "▲ Top Performers", font_size=14, bold=True, color=GREEN)

    _add_textbox(slide, Inches(0.6), Inches(3.7), Inches(12), Inches(0.35),
                 "▼ Needs Attention", font_size=14, bold=True, color=RED)

    # Top performer cards (hardcoded from channel_assessment data)
    top_cards = [
        ("6.39×", "Email",            "$527K spend",  "Most efficient channel — severely underfunded"),
        ("2.60×", "FB / Instagram",   "$1.65M spend", "#1 lead source — dual revenue & pipeline impact"),
        ("2.53×", "Paid Search",      "$2.1M spend",  "High-intent, above-blended-average ROAS"),
    ]

    bottom_cards = [
        ("1.61×", "Reddit",  "$588K spend",  "Lowest ROAS — skeptical audience, weak targeting"),
        ("1.81×", "Display", "$1.44M spend", "Broad programmatic — viewability & targeting issues"),
        ("1.83×", "TikTok",  "$1.21M spend", "Creative fatigue — needs refresh or budget cut"),
    ]

    card_w = Inches(4.0)
    card_h = Inches(2.1)
    card_xs = [Inches(0.6), Inches(4.8), Inches(9.0)]

    def _draw_channel_card(roas, channel, spend, desc, x, y, badge_color, strip_color):
        # Card background
        _add_rect(slide, x, y, card_w, card_h, LIGHT,
                  line_color=RGBColor(0xDD, 0xE3, 0xEC))
        # Left strip
        _add_rect(slide, x, y, Inches(0.08), card_h, strip_color)
        # ROAS badge
        _add_rect(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.9), Inches(0.42), badge_color)
        _add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.9), Inches(0.42),
                     roas, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Channel name
        _add_textbox(slide, x + Inches(1.2), y + Inches(0.2), Inches(2.6), Inches(0.45),
                     channel, font_size=15, bold=True, color=NAVY)
        # Spend
        _add_textbox(slide, x + Inches(1.2), y + Inches(0.65), Inches(2.6), Inches(0.3),
                     spend, font_size=11, color=GREY)
        # Description (wrapping)
        _add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(3.7), Inches(0.7),
                     desc, font_size=11, color=NAVY, wrap=True)

    # Draw top performer cards
    top_y = Inches(1.4)
    for i, (roas, channel, spend, desc) in enumerate(top_cards):
        _draw_channel_card(roas, channel, spend, desc, card_xs[i], top_y, TEAL, TEAL)

    # Draw bottom cards
    bottom_y = Inches(4.1)
    for i, (roas, channel, spend, desc) in enumerate(bottom_cards):
        _draw_channel_card(roas, channel, spend, desc, card_xs[i], bottom_y, RED, RED)


def _slide_budget_optimizer(prs: Presentation, report: dict):
    """Budget reallocation slide derived from MMM ROAS output."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                 "Budget Optimizer — Same Spend, More Revenue", font_size=24, bold=True, color=WHITE)

    # Subtitle
    _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.45),
                 "Shift $970K from low-ROAS channels into Email, Paid Search, and FB/IG. "
                 "Total budget unchanged at $15.0M.",
                 font_size=13, color=GREY)

    # Impact KPI boxes — equal widths, centered
    BOX_W = Inches(3.8)
    kpi_data = [
        (TEAL,  "$970K",        "Reallocated",       Inches(0.6)),
        (GREEN, "+$2.1M",       "Proj. Revenue Gain", Inches(4.7)),
        (NAVY,  "2.21×→2.35×", "Blended ROAS",  Inches(8.8)),
    ]
    for color, val, lbl, x in kpi_data:
        _add_rect(slide, x, Inches(1.85), BOX_W, Inches(1.3), color)
        _add_textbox(slide, x + Inches(0.1), Inches(1.85) + Inches(0.15),
                     BOX_W - Inches(0.2), Inches(0.5),
                     val, font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(1.85) + Inches(0.75),
                     BOX_W - Inches(0.2), Inches(0.4),
                     lbl, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Reallocation table
    rows = [
        ("Email",       "6.39×", "$527K",   "$1,327K", "+$800K",  "Scale 2.5×"),
        ("Paid Search", "2.53×", "$2,105K", "$2,205K", "+$100K",  "Increase"),
        ("FB / IG",     "2.60×", "$1,654K", "$1,724K", "+$70K",   "Increase"),
        ("TV/CTV",      "2.34×", "$5,606K", "$5,606K", "—",  "Hold (brand)"),
        ("Influencer",  "2.29×", "$1,854K", "$1,854K", "—",  "Hold"),
        ("TikTok",      "1.83×", "$1,209K", "$846K",   "−$363K", "Reduce 30%"),
        ("Display",     "1.81×", "$1,437K", "$1,006K", "−$431K", "Reduce 30%"),
        ("Reddit",      "1.61×", "$588K",   "$411K",   "−$176K", "Reduce 30%"),
    ]
    headers = ["Channel", "ROAS", "Current", "Optimized", "Change", "Action"]
    col_x   = [0.5, 2.4, 3.8, 5.2, 6.6, 8.0]
    col_w   = [1.7, 1.2, 1.2, 1.2, 1.2, 2.5]

    top = Inches(3.4)
    # Header row
    _add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.32), NAVY)
    for j, h in enumerate(headers):
        _add_textbox(slide, Inches(col_x[j]), top + Inches(0.04),
                     Inches(col_w[j]), Inches(0.28), h,
                     font_size=10, bold=True, color=WHITE)
    top += Inches(0.34)

    for r_idx, row in enumerate(rows):
        bg = RGBColor(0xF4, 0xF7, 0xFB) if r_idx % 2 == 0 else WHITE
        _add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.38), bg)
        for j, cell in enumerate(row):
            cell_color = NAVY
            if j == 4:  # Change column
                if cell.startswith("+"):
                    cell_color = GREEN
                elif cell.startswith("−"):
                    cell_color = RED
            elif j == 1:  # ROAS
                val = float(row[1].replace("×", ""))
                cell_color = GREEN if val >= 2.5 else (AMBER if val >= 1.9 else RED)
            _add_textbox(slide, Inches(col_x[j]), top + Inches(0.04),
                         Inches(col_w[j]), Inches(0.32), cell,
                         font_size=11, bold=(j == 4 and cell != "—"),
                         color=cell_color)
        top += Inches(0.40)

    _add_textbox(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.28),
                 "* Incremental revenue assumes 55% of gross ROAS delta applies at margin (accounts for diminishing returns as Email scales).",
                 font_size=9, color=GREY)


def _slide_funnel(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    # Navy header
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 "Sales Funnel Performance", font_size=24, bold=True, color=WHITE)

    # LEFT PANEL
    _add_textbox(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.35),
                 "Funnel Conversion", font_size=14, bold=True, color=NAVY)

    # Funnel stages
    stages = [
        ("Lead",        25019, 100),
        ("MQL",         9508,  38),
        ("SQL",         5754,  23),
        ("Opportunity", 3510,  14),
        ("Closed Won",  655,   2.6),
    ]
    stage_colors = [TEAL, TEAL, TEAL, AMBER, GREEN]

    for i, ((stage, count, pct), bar_color) in enumerate(zip(stages, stage_colors)):
        y = Inches(1.55) + i * Inches(0.95)
        # Stage label
        _add_textbox(slide, Inches(0.5), y, Inches(1.4), Inches(0.32),
                     stage, font_size=12, bold=True, color=NAVY)
        # Count
        _add_textbox(slide, Inches(2.0), y, Inches(1.2), Inches(0.32),
                     f"{count:,}", font_size=11, color=GREY, align=PP_ALIGN.RIGHT)
        # Gray track bar
        _add_rect(slide, Inches(3.4), y + Inches(0.05), Inches(3.0), Inches(0.22),
                  RGBColor(0xDD, 0xE3, 0xEC))
        # Fill bar (capped at full width)
        fill_w = Inches(3.0 * pct / 100)
        if fill_w > 0:
            _add_rect(slide, Inches(3.4), y + Inches(0.05), fill_w, Inches(0.22), bar_color)
        # Pct label
        _add_textbox(slide, Inches(3.4) + fill_w + Inches(0.05), y, Inches(0.5), Inches(0.28),
                     f"{pct}%" if isinstance(pct, int) else f"{pct:.1f}%",
                     font_size=10, color=TEAL)

    # Insight text below funnel
    _add_textbox(slide, Inches(0.5), Inches(6.3), Inches(6.0), Inches(0.45),
                 "2.62% win rate vs 5–8% CPG benchmark — 97.4% of leads are lost",
                 font_size=12, color=RED)

    # Vertical divider
    _add_rect(slide, Inches(6.9), Inches(1.0), Inches(0.02), Inches(5.8), GREY)

    # RIGHT PANEL
    funnel = report.get("insights", {}).get("funnel_health", {})
    recs = funnel.get("sales_recommendations", [])

    _add_textbox(slide, Inches(7.1), Inches(1.05), Inches(5.7), Inches(0.35),
                 "Key Recommendations", font_size=14, bold=True, color=NAVY)

    for i, rec in enumerate(recs[:4]):
        y = Inches(1.55) + i * Inches(1.3)
        # Small teal indicator
        _add_rect(slide, Inches(7.1), y + Inches(0.08), Inches(0.22), Inches(0.22), TEAL)
        # Rec text — first ~80 chars up to first em dash or period
        rec_text = rec
        for sep in [" —", ". ", ".\n"]:
            idx = rec.find(sep)
            if idx != -1 and idx < 90:
                rec_text = rec[:idx]
                break
        if len(rec_text) > 90:
            rec_text = rec_text[:90] + "…"
        _add_textbox(slide, Inches(7.45), y, Inches(5.3), Inches(0.55),
                     rec_text, font_size=12, color=NAVY, wrap=True)


def _slide_risks_opportunities(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    # Navy header
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 "Risks & Opportunities", font_size=24, bold=True, color=WHITE)

    insights = report.get("insights", {})
    risks = insights.get("risks", [])
    opps  = insights.get("opportunities", [])

    sev_color = {"high": RED, "medium": AMBER, "low": GREEN}

    # LEFT COLUMN — Risks
    _add_textbox(slide, Inches(0.5), Inches(1.02), Inches(5.5), Inches(0.4),
                 "⚠ Risks", font_size=15, bold=True, color=RED)

    for i, risk in enumerate(risks[:3]):
        y = Inches(1.55) + i * Inches(1.55)
        sev = risk.get("severity", "medium")
        border_color = sev_color.get(sev, AMBER)

        # Left border rect
        _add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(1.3), border_color)

        # Risk title — first sentence only (up to first period or 90 chars)
        risk_text = risk.get("risk", "")
        period_idx = risk_text.find(". ")
        if period_idx != -1 and period_idx < 90:
            risk_title = risk_text[:period_idx]
        elif len(risk_text) > 90:
            risk_title = risk_text[:90] + "…"
        else:
            risk_title = risk_text

        _add_textbox(slide, Inches(0.68), y, Inches(5.6), Inches(0.55),
                     risk_title, font_size=12, bold=True, color=NAVY, wrap=True)

        # Action text
        action_text = risk.get("action", "")
        if len(action_text) > 80:
            action_text = action_text[:80] + "…"
        _add_textbox(slide, Inches(0.68), y + Inches(0.58), Inches(5.6), Inches(0.45),
                     "→ " + action_text, font_size=11, color=GREY, wrap=True)

    # Vertical divider
    _add_rect(slide, Inches(6.65), Inches(1.0), Inches(0.02), Inches(6.0), GREY)

    # RIGHT COLUMN — Opportunities
    _add_textbox(slide, Inches(6.9), Inches(1.02), Inches(5.9), Inches(0.4),
                 "✦ Opportunities", font_size=15, bold=True, color=GREEN)

    for i, opp in enumerate(opps[:3]):
        y = Inches(1.55) + i * Inches(1.55)

        # Green left border
        _add_rect(slide, Inches(6.9), y, Inches(0.06), Inches(1.3), GREEN)

        # Opp title — first sentence (up to em dash or 90 chars)
        opp_text = opp.get("opportunity", "")
        for sep in [" —", ". "]:
            idx = opp_text.find(sep)
            if idx != -1 and idx < 90:
                opp_title = opp_text[:idx]
                break
        else:
            opp_title = opp_text[:90] + "…" if len(opp_text) > 90 else opp_text

        _add_textbox(slide, Inches(7.08), y, Inches(5.7), Inches(0.55),
                     opp_title, font_size=12, bold=True, color=NAVY, wrap=True)

        # Impact + timeframe
        impact = opp.get("estimated_impact", "")
        timeframe = opp.get("timeframe", "")
        _add_textbox(slide, Inches(7.08), y + Inches(0.58), Inches(5.7), Inches(0.45),
                     f"Impact: {impact}  ·  {timeframe}",
                     font_size=11, color=GREY, wrap=True)


def _slide_actions(prs: Presentation, report: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)

    # Teal header
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), TEAL)
    _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
                 "Recommended Actions", font_size=24, bold=True, color=WHITE)

    actions = report.get("insights", {}).get("recommended_actions", [])
    urgency_color = {"Immediate": RED, "This Quarter": AMBER, "Strategic": TEAL}
    urgency_label = {"Immediate": "IMMEDIATE", "This Quarter": "THIS QTR", "Strategic": "STRATEGIC"}

    for i, action in enumerate(actions[:5]):
        y = Inches(1.0) + i * Inches(1.12)
        urgency = action.get("urgency", "This Quarter")
        badge_color = urgency_color.get(urgency, TEAL)
        badge_text  = urgency_label.get(urgency, urgency.upper()[:9])

        # Number circle (navy rect)
        _add_rect(slide, Inches(0.5), y + Inches(0.18), Inches(0.42), Inches(0.42), NAVY)
        _add_textbox(slide, Inches(0.5) + Inches(0.1), y + Inches(0.18) + Inches(0.05),
                     Inches(0.42) - Inches(0.1), Inches(0.32),
                     str(i + 1), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Urgency badge (top-right)
        _add_rect(slide, Inches(12.0), y + Inches(0.18), Inches(1.0), Inches(0.3), badge_color)
        _add_textbox(slide, Inches(12.0), y + Inches(0.18), Inches(1.0), Inches(0.3),
                     badge_text, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Action title (full text, wide box, word wrap)
        action_text = action.get("action", "")
        if len(action_text) > 130:
            action_text = action_text[:130] + "…"
        _add_textbox(slide, Inches(1.1), y + Inches(0.1), Inches(10.7), Inches(0.48),
                     action_text, font_size=13, bold=True, color=NAVY, wrap=True)

        # Impact text
        _add_textbox(slide, Inches(1.1), y + Inches(0.6), Inches(10.7), Inches(0.35),
                     "↗ " + action.get("expected_impact", ""),
                     font_size=10, color=GREY, wrap=True)

        # Thin divider below item
        if i < 4:
            _add_rect(slide, Inches(0.5), y + Inches(1.1), Inches(12.3), Inches(0.01), GREY)


# ── Main entry point ──────────────────────────────────────────────────────────

def generate_deck(
    insights_path: str | Path | None = None,
    output_path: str | Path | None = None,
) -> Path:
    """
    Generate an executive PPTX deck from an insights report JSON.

    Args:
        insights_path: Path to insights_report.json (defaults to outputs/insights_report.json).
        output_path:   Where to save the deck (defaults to outputs/executive_deck.pptx).

    Returns:
        Path to the generated PPTX file.
    """
    insights_path = Path(insights_path or OUTPUT_DIR / "insights_report.json")
    output_path   = Path(output_path   or OUTPUT_DIR / "executive_deck.pptx")

    if not insights_path.exists():
        raise FileNotFoundError(
            f"Insights report not found at {insights_path}. "
            "Run `python insights/generate_insights.py` first."
        )

    report = json.loads(insights_path.read_text())

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs, report)
    _slide_exec_summary(prs, report)
    _slide_channel_highlights(prs, report)
    _slide_budget_optimizer(prs, report)
    _slide_funnel(prs, report)
    _slide_risks_opportunities(prs, report)
    _slide_actions(prs, report)

    OUTPUT_DIR.mkdir(exist_ok=True)
    prs.save(output_path)
    print(f"Deck saved to {output_path}  ({prs.slides.__len__()} slides)")
    return output_path


def upload_to_google_slides(
    pptx_path: Path,
    folder_id: str | None = None,
) -> str:
    """
    Upload a PPTX file to Google Drive and convert it to Google Slides.

    Authentication uses the service account at GOOGLE_APPLICATION_CREDENTIALS,
    or falls back to OAuth via a browser prompt if that is not set.

    Args:
        pptx_path:  Path to the .pptx file to upload.
        folder_id:  Optional Google Drive folder ID to place the file in.

    Returns:
        The URL of the created Google Slides presentation.
    """
    import sys
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload
    from google.oauth2 import service_account
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request
    import google.auth
    import pickle

    SCOPES = ["https://www.googleapis.com/auth/drive.file"]
    creds = None

    sa_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS", "")
    if sa_path and Path(sa_path).exists():
        # Service account auth (production / CI)
        creds = service_account.Credentials.from_service_account_file(
            sa_path, scopes=SCOPES
        )
        print("  Authenticating via service account...")
    else:
        # OAuth flow (local dev) — caches token in outputs/.google_token.pkl
        token_path = OUTPUT_DIR / ".google_token.pkl"
        if token_path.exists():
            with open(token_path, "rb") as f:
                creds = pickle.load(f)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                # Requires a client_secrets.json — see docs/google_oauth_setup.md
                secrets_path = Path(__file__).resolve().parent.parent / "client_secrets.json"
                if not secrets_path.exists():
                    raise FileNotFoundError(
                        "Google OAuth credentials not found.\n"
                        "Either set GOOGLE_APPLICATION_CREDENTIALS to a service account JSON,\n"
                        f"or place a client_secrets.json at {secrets_path}.\n"
                        "See docs/google_oauth_setup.md for setup instructions."
                    )
                flow = InstalledAppFlow.from_client_secrets_file(str(secrets_path), SCOPES)
                creds = flow.run_local_server(port=0)

            with open(token_path, "wb") as f:
                pickle.dump(creds, f)
            print("  OAuth token cached for future runs.")

    service = build("drive", "v3", credentials=creds)

    # File metadata — converting to Google Slides on upload
    file_name = pptx_path.stem  # filename without .pptx
    metadata = {
        "name": file_name,
        "mimeType": "application/vnd.google-apps.presentation",
    }
    if folder_id:
        metadata["parents"] = [folder_id]

    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )

    print(f"  Uploading '{file_name}' to Google Slides...")
    file = (
        service.files()
        .create(body=metadata, media_body=media, fields="id,webViewLink")
        .execute()
    )

    url = file.get("webViewLink", f"https://docs.google.com/presentation/d/{file['id']}/edit")
    print(f"  Google Slides URL: {url}")
    return url


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate executive deck from insights report.")
    parser.add_argument(
        "--upload",
        action="store_true",
        help="Upload the generated PPTX to Google Slides after saving.",
    )
    parser.add_argument(
        "--folder-id",
        default=None,
        help="Google Drive folder ID to upload into (optional).",
    )
    args = parser.parse_args()

    deck_path = generate_deck(output_path="/Users/jessleung/Downloads/CPG_Retail_Platform_Showcase_v2.pptx")

    if args.upload:
        try:
            url = upload_to_google_slides(deck_path, folder_id=args.folder_id)
            print(f"\nPresentation ready: {url}")
        except FileNotFoundError as e:
            print(f"\nUpload failed: {e}")
        except Exception as e:
            print(f"\nUpload failed: {e}")
            raise
